# SPDX-License-Interier: AGPL-3.0-or-later
# Copyright (c) 2025 Robert Čmrlec

from pptx import Presentation  

# load the presentation
prs = Presentation("ZGO_1.pptx")

# iterate through slides and shapes
for i, slide in enumerate(prs.slides, start=1):
    print(f"--- Slide {i} ---")
    for shape in slide.shapes: 
        if hasattr(shape, "text"): 
            print(shape.text)
